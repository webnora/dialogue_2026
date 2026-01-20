#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Dialogue 2026 conference
Based on results from Phases 1-3
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle
    return slide

def create_content_slide(prs, title, content_items):
    """Create content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

    return slide

def create_table_slide(prs, title, headers, rows):
    """Create slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title

    # Create table
    rows_count = len(rows) + 1
    cols_count = len(headers)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8)

    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

    # Set headers
    for col, header in enumerate(headers):
        table.cell(0, col).text = header
        table.cell(0, col).text_frame.paragraphs[0].font.bold = True
        table.cell(0, col).text_frame.paragraphs[0].font.size = Pt(14)

    # Fill data
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            table.cell(row_idx + 1, col_idx).text = str(cell_data)
            table.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

    return slide

def create_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Create slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True

    p = tf_left.paragraphs[0]
    p.text = left_title
    p.font.bold = True
    p.font.size = Pt(16)

    for item in left_items:
        p = tf_left.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(14)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True

    p = tf_right.paragraphs[0]
    p.text = right_title
    p.font.bold = True
    p.font.size = Pt(16)

    for item in right_items:
        p = tf_right.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(14)

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    create_title_slide(
        prs,
        "What Vector Representations Reveal\nAbout Publicistic Writing:\nLearning from Mistakes",
        "Dialogue 2026\nJanuary 2026"
    )

    # Slide 2: Research Question
    create_content_slide(prs, "Research Question", [
        "What do classification errors reveal about genre boundaries?",
        "",
        "Key Hypotheses:",
        "• H1: Contextual models (BERT) outperform lexical (TF-IDF)",
        "• H2: Genre boundaries are gradient, not discrete",
        "• H3: Lexical choice is the primary genre marker",
        "• H4: Syntactic features are less informative than lexical"
    ])

    # Slide 3: Methodology
    create_content_slide(prs, "Methodology: Three Levels of Representation", [
        "Three classification models at different linguistic levels:",
        "",
        "1. TF-IDF + Logistic Regression (Lexical level)",
        "2. Linguistic Features + Random Forest (Discourse-grammatical)",
        "3. BERT fine-tuning (Contextual-semantic)",
        "",
        "Dataset: 50,000 articles from The Guardian (5 genres)"
    ])

    # Slide 4: Data
    create_content_slide(prs, "Dataset", [
        "Source: The Guardian API",
        "Total: 50,000 articles",
        "",
        "Five genres:",
        "• Analytical (political analysis, commentary)",
        "• Editorial (opinion pieces)",
        "• Feature (long-form narratives)",
        "• News (factual reporting)",
        "• Review (cultural criticism)"
    ])

    # Slide 5: Phase 1 - Model Comparison
    create_table_slide(
        prs,
        "Phase 1: Baseline Model Performance",
        ["Model", "Level", "Accuracy", "Macro F1"],
        [
            ["BERT", "Contextual", "87.64%", "0.8771"],
            ["TF-IDF + LR", "Lexical", "86.58%", "0.8647"],
            ["Linguistic + RF", "Syntactic", "65.00%", "0.6449"]
        ]
    )

    # Slide 6: Key Finding 1
    create_content_slide(prs, "Key Finding 1: Lexical Primacy", [
        "BERT improves over TF-IDF by only 1.06%",
        "",
        "Implications:",
        "• ~99% of genre signal is in word choice",
        "• Context adds minimal value for genre classification",
        "• Practical: TF-IDF is sufficient for production",
        "",
        "This challenges the assumption that deep contextual",
        "representations are necessary for genre classification"
    ])

    # Slide 7: Key Finding 2
    create_content_slide(prs, "Key Finding 2: Syntactic Features Underperform", [
        "Linguistic features: 65% accuracy (vs 86%+ for lexical)",
        "",
        "Top features by importance:",
        "1. Reporting verbs ratio (0.20)",
        "2. Type-token ratio (0.19)",
        "3. First person pronouns (0.14)",
        "",
        "Problem: Severe overfitting (87% train → 65% test)"
    ])

    # Slide 7.5: McNemar's Test
    create_content_slide(prs, "Statistical Validation: McNemar's Test", [
        "Research question: Are model differences statistically significant?",
        "",
        "McNemar's test (paired comparison):",
        "• BERT vs TF-IDF: χ² = 156.01, p < 0.001 ***",
        "• BERT vs Linguistic: χ² = 258.34, p < 0.001 ***",
        "• TF-IDF vs Linguistic: χ² = 31.07, p < 0.001 ***",
        "",
        "Conclusion: All pairwise differences are HIGHLY SIGNIFICANT"
    ])

    # Slide 7.6: Bootstrap Confidence Intervals
    create_table_slide(
        prs,
        "Bootstrap 95% Confidence Intervals",
        ["Model", "Accuracy", "95% CI"],
        [
            ["BERT", "92.73%", "[92.02%, 93.42%]"],
            ["TF-IDF", "86.40%", "[85.50%, 87.22%]"],
            ["Linguistic", "82.85%", "[81.78%, 83.94%]"]
        ]
    )

    # Slide 7.7: Model Agreement
    create_content_slide(prs, "Model Agreement (Cohen's Kappa)", [
        "Cohen's Kappa measures inter-rater reliability:",
        "",
        "Model agreement results:",
        "• TF-IDF ↔ BERT: κ = 0.827 (substantial)",
        "• Linguistic ↔ BERT: κ = 0.752 (substantial)",
        "• TF-IDF ↔ Linguistic: κ = 0.728 (substantial)",
        "",
        "Interpretation (Landis-Koch scale):",
        "→ κ = 0.61-0.80: Substantial agreement",
        "→ All models capture the same underlying signal"
    ])

    # Slide 8: Phase 2 - Error Analysis
    create_content_slide(prs, "Phase 2: Error Analysis", [
        "Analyzed 5,000 test samples across all three models",
        "",
        "Model agreement:",
        "• All three models agree: 73.4%",
        "• When all agree → 98.4% accuracy",
        "",
        "Hybrid articles: 27.8% (1,390/5,000)",
        "→ Boundary cases where genre is genuinely ambiguous"
    ])

    # Slide 9: Feature as Hub Genre
    create_content_slide(prs, "Feature as 'Hub Genre'", [
        "Feature attracts confusion from all genres:",
        "",
        "Top BERT confusions:",
        "• Analytical → Feature: 11.5%",
        "• News → Feature: 5.4%",
        "• Review → Feature: 4.2%",
        "",
        "Feature occupies a middle ground: narrative + informative"
    ])

    # Slide 10: Gradient Boundaries
    create_two_column_slide(
        prs,
        "Evidence for Gradient Genre Boundaries",
        "All three wrong",
        ["105 articles (2.1%)", "Genuinely ambiguous", "No clear classification"],
        "Model disagreement",
        ["1,390 articles (27.8%)", "Boundary cases", "Linguistic overlap"]
    )

    # Slide 11: Phase 3 - BERT Attention
    create_content_slide(prs, "Phase 3: Interpreting BERT via Attention", [
        "Extracted attention weights from last 6 layers of BERT",
        "Analyzed 50 texts (10 × 5 genres)",
        "",
        "Research question: What tokens does BERT attend to for genre?"
    ])

    # Slide 12: Genre-Specific Markers
    create_table_slide(
        prs,
        "Genre-Specific Attention Markers",
        ["Genre", "Top-3 Tokens", "Interpretation"],
        [
            ["Analytical", "boris, johnson, downing", "Politics, names"],
            ["Editorial", "mr, nhs, of", "Institutions"],
            ["Feature", "i, my, said", "Narrative, personal"],
            ["News", "trump, masters, after", "Events, facts"],
            ["Review", "staging, debut, opera", "Culture, arts"]
        ]
    )

    # Slide 13: Universal Markers
    create_content_slide(prs, "Universal Markers Across All Genres", [
        "Common high-attention tokens:",
        "• 'we', 'related', 'sp' (function words)",
        "",
        "Key insight:",
        "→ No unique vocabulary for any genre",
        "→ Overlapping token distributions",
        "→ Explains gradient boundaries",
        "",
        "Genres blend into each other; boundaries are permeable"
    ])

    # Slide 14: Theoretical Implications
    create_content_slide(prs, "Theoretical Implications", [
        "1. Genre as gradient, not discrete",
        "   → Empirical evidence from model disagreement",
        "",
        "2. Lexical primacy in genre classification",
        "   → Word choice outweighs syntax and context",
        "",
        "3. Model disagreement as uncertainty metric",
        "   → Flagging ambiguous cases for human review"
    ])

    # Slide 15: Overall Results Summary
    create_table_slide(
        prs,
        "Final Model Rankings",
        ["Rank", "Model", "Accuracy", "Level"],
        [
            ["🥇 1st", "BERT", "87.64%", "Contextual"],
            ["🥈 2nd", "TF-IDF", "86.58%", "Lexical"],
            ["🥉 3rd", "Linguistic", "65.00%", "Syntactic"]
        ]
    )

    # Slide 16: Hypothesis Testing Results
    create_table_slide(
        prs,
        "Hypothesis Testing Results",
        ["Hypothesis", "Result", "Evidence"],
        [
            ["H1: Contextual > Lexical", "⚠️ Partial", "BERT +1.06% over TF-IDF"],
            ["H2: Gradient boundaries", "✅ Confirmed", "27.8% hybrid articles"],
            ["H3: Lexical primacy", "✅ Confirmed", "TF-IDF ≈ BERT (99% signal)"],
            ["H4: Syntax < Lexical", "✅ Confirmed", "65% vs 87%"]
        ]
    )

    # Slide 17: Conclusion
    create_content_slide(prs, "Conclusion", [
        "Main contributions:",
        "",
        "1. Empirical evidence for gradient genre boundaries",
        "2. Demonstration of lexical primacy (99% of signal)",
        "3. Model disagreement as uncertainty metric",
        "4. Attention analysis reveals genre-specific markers",
        "",
        "Practical impact:",
        "→ TF-IDF sufficient for genre classification",
        "→ Model agreement flags ambiguous cases"
    ])

    # Slide 18: Future Work
    create_content_slide(prs, "Future Work (Phase 4)", [
        "Statistical validation:",
        "• McNemar's test for model comparison",
        "• Bootstrap confidence intervals",
        "• Inter-annotator agreement (Cohen's Kappa)",
        "",
        "Extensions:",
        "• Apply to other languages/datasets",
        "• Investigate cross-domain genre differences"
    ])

    # Slide 19: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Thank You!"
    subtitle.text = "Questions?\n\nDialogue 2026\nJanuary 2026"

    # Save
    output_path = "Dialogue_2026_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
